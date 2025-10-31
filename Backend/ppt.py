import io
import json
import os

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from urllib.parse import quote_plus
from dotenv import load_dotenv
from dotenv import get_key
from rich import print
import requests
import base64


class LLM:
    USER = "user"
    ASSISTANT = "assistant"
    SYSTEM = "system"
    def __init__(
            self,
            messages: list[dict[str, str]] = [],
            model: str = "rohan/tune-gpt-4o",
            temperature: float = 0.0,
            system_prompt: str = "",
            max_tokens: int = 4092,
            verbose: bool = False,
            api_key: str | None = None
    ) -> None:
        """
        Initializes the LLM with the given parameters.

        Args
        ----
        messages: list[dict[str, str]]
            A list of messages to initialize the LLM with.
        model: str
            The model to use for the LLM.
        temperature: float
            The temperature to use for the LLM.
        system_prompt: str
            The system prompt to use for the LLM.
        max_tokens: int
            The maximum number of tokens to use for the LLM.
        verbose: bool
            Whether to print the response from the LLM.
        api_key: str | None
            The API key to use for the LLM.
        
        example:
        >>> llm = LLM()
        >>> llm.add_message("user", "Hello, how are you?")
        >>> llm.add_message("assistant", "I'm doing well, thank you!")
        >>> llm.run("Hello, how are you?")
        >>> "I'm doing well, thank you!"
        """
        self.api_key = api_key if api_key else get_key(".env","TUNE_STUDIO_API_KEY")
        self.session =  requests.session()
        self.messages = messages
        self.model = model
        self.temperature = temperature
        self.system_prompt = system_prompt
        self.max_tokens = max_tokens
        self.verbose = verbose
        "" if not system_prompt else self.add_message(self.SYSTEM, system_prompt)
    def run(self, prompt: str|None = None) -> str:
        """
        Runs the LLM with the given prompt.

        Args
        ----
        prompt: str
            The prompt to use for the LLM.

        Returns
        -------
        str
            The response from the LLM.

        example:
        >>> llm = LLM()
        >>> llm.add_message("user", "Hello, how are you?")
        >>> llm.add_message("assistant", "I'm doing well, thank you!")
        >>> llm.run("Hello, how are you?")
        """

        "" if not prompt else self.add_message("user", prompt)
        url = "https://proxy.tune.app/chat/completions"
        headers = {
            "Authorization": self.api_key,
            "Content-Type": "application/json",
        }
        data = {
        "temperature": self.temperature,
        
            "messages":  self.messages,
            "model": self.model,
            "stream": False,
            "frequency_penalty":  0.0,
            "max_tokens": self.max_tokens
        }
        response = self.session.post(url, headers=headers, json=data)
        "" if not prompt else self.messages.pop()
        return response.json()["choices"][0]["message"]["content"]

    def add_message(self, role: str, content: str, base64_image: str = "") -> None:
        """
        Adds a message to the LLM with the given role and content.

        Args
        ----
        role: str
            The role of the message.
        content: str
            The content of the message.
        base64_image: str
            The base64 image of the message.

        example
        -------
        >>> llm = LLM()
        >>> llm.add_message("user", "Hello, how are you?")
        >>> llm.add_message("assistant", "I'm doing well, thank you!")
        >>> llm.run("Hello, how are you?")
        >>> "I'm doing well, thank you!"
        """
        if content and base64_image:
            self.messages.append({
                "role": role,
                "content": [
                {
                    "type": "text",
                    "text": content
                },
                {
                    "type": "image_url", 
                    "image_url": 
                {
                    "url": f"data:image/png;base64,{base64_image}"
                }
                }
                ]
                })
        elif base64_image:
            self.messages.append(
                {
                "role": role, 
                "content": [
                    {
                    "type": "image_url", 
                    "image_url": 
                        {
                        "url": f"data:image/png;base64,{base64_image}"
                        }
                    }
                ]       
                }
                )
        elif content:
            self.messages.append(
                {
                "role": role, 
                "content": [
                    {
                    "type": "text",
                    "text": content
                }
                ]       
                }
                )
        else:
            raise ValueError("Both content and base64_image are None")

    
    def __getitem__(self, index) -> dict[str, str] | list[dict[str, str]]:
        """
        Returns the message at the given index.

        Args
        ----
        index: int
            The index of the message to return.

        Returns
        -------
        dict[str, str] | list[dict[str, str]]
            The message at the given index.

        example
        -------
        >>> llm = LLM()
        >>> llm.add_message("user", "Hello, how are you?")
        >>> llm.add_message("assistant", "I'm doing well, thank you!")
        >>> llm.run("Hello, how are you?")
        >>> "I'm doing well, thank you!"
        >>> llm[1]
        >>> llm[1:]
        """
        if isinstance(index, slice):
            return self.messages[index]
        elif isinstance(index, int):
            return self.messages[index]
        else:
            raise TypeError("Invalid argument type")

    def __setitem__(self, index, value) -> None:
        """
        Sets the message at the given index to the given value.

        Args
        ----
        index: int
            The index of the message to set.
        value: dict[str, str] | list[dict[str, str]]
            The value to set the message to.

        example
        -------
        >>> llm = LLM()
        >>> llm.add_message("user", "Hello, how are you?")
        >>> llm.add_message("assistant", "I'm doing well, thank you!")
        >>> llm.run("Hello, how are you?")
        >>> "I'm doing well, thank you!"
        >>> llm[1] = "I'm doing well, thank you!"
        """

        if isinstance(index, slice):
            self.messages[index] = value
        elif isinstance(index, int):
            self.messages[index] = value
        else:
            raise TypeError("Invalid argument type")

def FileToBase64(file_path:str):
    """
    Convert image file to base64 string.

    Args
    ----
    file_path : str

    Returns
    -------
    base64_image : str
    """
    with open(file_path, "rb") as image_file:
        encoded_image = base64.b64encode(image_file.read()).decode("utf-8")
    return encoded_image


dir_path = 'presentations'

load_dotenv()
API_KEY = os.getenv('PEXELS_API_KEY')

def parse_response(response:str):
    slides = response.split('\n\n')
    slides_content = []
    for slide in slides:
        lines = slide.split('\n')
        title_line = lines[0]
        if ': ' in title_line:
            title = title_line.split(': ', 1)[1]  # Extract the title after 'Slide X: '
        else:
            title = title_line
        content_lines = [line for line in lines[1:] if line != 'Content:']  # Skip line if it is 'Content:'
        content = '\n'.join(content_lines)  # Join the lines to form the content
        # Extract the keyword from the line that starts with 'Keyword:'
        keyword_line = [line for line in lines if 'Keyword:' or 'Keywords:' in line][0]
        keyword = keyword_line.split(': ', 1)[1]
        slides_content.append({'title': title, 'content': content, 'keyword': keyword})
    return slides_content

def search_pexels_images(keyword, page=1, per_page=1, order_by='relevant', color=None, orientation=None):
    access_key = os.getenv("UNSPLASH_ACCESS_KEY")
    if not access_key:
        return "Access key not found."

    # Unsplash API URL for searching photos
    url = "https://api.unsplash.com/search/photos"
    
    # Prepare headers and parameters for the request
    headers = {
        "Authorization": f"Client-ID {access_key}"
    }
    
    params = {
        "query": keyword,
        "page": page,
        "per_page": per_page,
        "order_by": order_by
    }
    
    # Optional parameters
    if color:
        params["color"] = color
    if orientation:
        params["orientation"] = orientation
    
    # Make the GET request to Unsplash API
    response = requests.get(url, headers=headers, params=params)
    
    if response.status_code == 200:
        data = response.json()
        if data['results']:
            # Return URLs of the images
            return [result['urls']['regular'] for result in data['results']][0]
        else:
            return "No images found for the keyword."
    else:
        return f"Error: {response.status_code} - {response.text}"


def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(slides_content: list[str], template_choice, presentation_title, presenter_name, insert_image):
    template_path = os.path.join(dir_path, f"{template_choice}.pptx")

    prs = Presentation(template_path)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title

    #add subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"

    if template_choice == 'dark_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color

    elif template_choice == 'bright_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Arial'
                run.font.color.rgb = RGBColor(255, 20, 147)  # RGB for deep pink color

    # add content slides
    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = slide_content['content']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for white color

        if insert_image:
            # fetch image URL from Pixabay based on the slide's title
            image_url = search_pexels_images(slide_content['keyword'])
            print("Image URL:", image_url) #debug
            if image_url is not None:
                # download the image
                image_data = requests.get(image_url).content
                # load image into BytesIO object
                image_stream = io.BytesIO(image_data)
                # add the image at the specified position
                slide_width = Inches(20)
                slide_height = Inches(15)

                image_width = Inches(8)  # width of image
                image_height = Inches(5)  # height of image

                left = slide_width - image_width  # calculate left position
                top = slide_height - image_height - Inches(4)  # calculate top position

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

    # add credits slide
    slide = prs.slides.add_slide(content_slide_layout)
    if template_choice == 'dark_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 165, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 255, 255)

    elif template_choice == 'bright_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(255, 20, 147)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    else:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    # Save the presentation
    prs.save(os.path.join('generated', 'generated_presentation.pptx'))



def main(topic:str, number_of_slides:int, theme:str, presentation_title, presenter_name):
    m = [
    {
        'role': 'system',
        'content': '''You are an assistant that gives the idea for PowerPoint presentations. When answering, give the user the  
summarized content for each slide based on the number of slide. And the format of the answer must be Slide X(the number of theslide): {title of the content} /n Content: /n content with some bullet points.Keyword: /n Give the most important 
keyword(within two words) that represents the slide for each one\nexample:Slide 1: Introduction to Python\nContent:\n- Brief  
overview of Python programming language\n- History and creator of Python (Guido van Rossum)\n- High-level, interpreted        
language\n- Used for web development, automation, data analysis, etc.\n\nKeyword: Python History\n\n-     
Comparison with other popular languages'''
    },
    {
        'role': 'user',
        'content': f'''I want you to come up with the idea for the PowerPoint. The number of slides is {number_of_slides}. The content is: {topic}.The title of content for each slide must be unique, and extract the most important keyword within two words for each slide. Summarize the content for each slide.'''
    }
]   
    for i in range(5):
        try:
            resp = LLM(messages=m).run()
            parsed = parse_response(resp)
            break
        except Exception as e:
            from rich import print
            print(e)
            continue
        if i == 4:
            raise Exception("Failed to generate ppt")
    create_ppt(parsed,
               theme,
               presenter_name,
               presentation_title,
               insert_image=True
               )
    
    
    

def PPT(topic:str):
    choice = ["bright_modern","dark_modern","simple"]
    main(topic, 4, choice[1], "Title", "vinay")
    os.startfile(os.path.join('generated', 'generated_presentation.pptx'))


if __name__ == '__main__':
    PPT("java")